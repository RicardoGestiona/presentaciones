"""
build_pptx_native.py — PPTX 100 % nativo editable + PDF.

Todos los elementos son objetos PowerPoint (TextBox, Shape, Table).
Sin imágenes de fondo. Fidelidad visual máxima posible con python-pptx.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Optional

from bs4 import BeautifulSoup, Tag
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

# ── rutas ──────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parent.parent
INDEX    = ROOT / "index.html"
OUT_PPTX = ROOT / "certificacion-inap.pptx"
OUT_PDF  = ROOT / "certificacion-inap.pdf"

# ── dimensiones 1280×720 ───────────────────────────────────────────────────
W, H       = 1280, 720
EMU        = 9525          # 1 CSS px = 9525 EMU (96 DPI)
W_EMU      = W * EMU
H_EMU      = H * EMU

def px(n: float) -> Emu:
    return Emu(int(n * EMU))

# ── paleta gestiona ────────────────────────────────────────────────────────
C = {
    "primary":      RGBColor(0,   109, 133),   # #006d85
    "primary_dark": RGBColor(0,    77,  94),   # #004d5e
    "primary_l":    RGBColor(0,   223, 178),   # #00dfb2
    "accent":       RGBColor(95,  255, 223),   # #5fffdf
    "bg_dark":      RGBColor(0,    48,  64),   # #003040
    "bg_light":     RGBColor(236, 240, 243),   # #ecf0f3
    "white":        RGBColor(255, 255, 255),
    "text_dark":    RGBColor(0,    48,  64),   # #003040
    "text_body":    RGBColor(26,   58,  68),   # #1a3a44
    "text_muted":   RGBColor(61,   92, 102),   # #3d5c66
    "border":       RGBColor(212, 222, 227),   # #d4dee8
}
FONT = "Roboto"

# ── helpers de bajo nivel ──────────────────────────────────────────────────

def solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape) -> None:
    shape.line.fill.background()

def rect(slide, x, y, w, h, color: RGBColor, *, rounded=False):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if rounded:
        s = slide.shapes.add_shape(5, px(x), px(y), px(w), px(h))  # 5 = rounded rect
        s.adjustments[0] = 0.04   # radio pequeño ≈ 12–14 px
    else:
        s = slide.shapes.add_shape(1, px(x), px(y), px(w), px(h))
    solid(s, color)
    no_line(s)
    return s

def textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(px(x), px(y), px(w), px(h))

def para(tf, text: str, pt: float, bold=False, color=None, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0, word_wrap=True, italic=False):
    p   = tf.add_paragraph()
    p.alignment    = align
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    run = p.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color or C["text_dark"]
    return p

def para_first(tf, text: str, pt: float, bold=False, color=None,
               align=PP_ALIGN.LEFT, space_before=0, space_after=0):
    """Usa el primer párrafo (ya existe al crear el tf)."""
    p   = tf.paragraphs[0]
    p.alignment    = align
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    run = p.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.color.rgb = color or C["text_dark"]
    return p

def add_run(paragraph, text: str, pt: float, bold=False, color=None, italic=False):
    run = paragraph.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color or C["text_dark"]
    return run

# ── tabla helpers ──────────────────────────────────────────────────────────

def cell_solid(cell, color: RGBColor) -> None:
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    sf   = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(sf,   qn("a:srgbClr"))
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")

def cell_no_border(cell) -> None:
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for b in ("a:lnL","a:lnR","a:lnT","a:lnB"):
        el = tcPr.find(qn(b))
        if el is None:
            el = etree.SubElement(tcPr, qn(b))
        for ch in list(el): el.remove(ch)
        etree.SubElement(el, qn("a:noFill"))

def cell_write(cell, text: str, pt: float, bold=False, color=None,
               align=PP_ALIGN.LEFT, margin_top=6, margin_left=8) -> None:
    tf = cell.text_frame
    tf.clear()
    tf.margin_top    = Pt(margin_top)
    tf.margin_bottom = Pt(2)
    tf.margin_left   = Pt(margin_left)
    tf.margin_right  = Pt(4)
    tf.word_wrap     = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.color.rgb = color or C["white"]

# ── componentes comunes ────────────────────────────────────────────────────

def slide_bg(slide) -> None:
    """Fondo blanco."""
    r = rect(slide, 0, 0, W, H, C["white"])

def slide_top_bar(slide) -> None:
    """Barra de color de 6px en la parte superior."""
    rect(slide, 0, 0, W, 6, C["primary"])

def slide_title_text(slide, text: str, y=18, w_offset=0) -> None:
    tb = textbox(slide, 40, y, W - 80 - w_offset, 44)
    tf = tb.text_frame
    tf.word_wrap = False
    para_first(tf, text, 20, bold=True, color=C["primary_dark"])

def slide_number_label(slide, n: int) -> None:
    tb = textbox(slide, W - 48, H - 26, 40, 20)
    tf = tb.text_frame
    para_first(tf, str(n), 10, color=C["text_muted"], align=PP_ALIGN.RIGHT)

def divider(slide, x=40, y=64, w=None) -> None:
    w = w or (W - 80)
    rect(slide, x, y, w, 2, C["accent"])

# ── bullet list ────────────────────────────────────────────────────────────

def bullet_list(slide, items: list[dict], x, y, w, h, pt=11.5):
    """
    items: [{"text": str, "bold_prefix": str|None}]
    """
    tb = textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after  = Pt(3)

        # bullet ▸
        r0 = p.add_run()
        r0.text           = "▸  "
        r0.font.name      = FONT
        r0.font.size      = Pt(pt)
        r0.font.bold      = True
        r0.font.color.rgb = C["primary"]

        if item.get("bold_prefix"):
            rb = p.add_run()
            rb.text           = item["bold_prefix"]
            rb.font.name      = FONT
            rb.font.size      = Pt(pt)
            rb.font.bold      = True
            rb.font.color.rgb = C["text_dark"]
            rest = item["text"][len(item["bold_prefix"]):]
            if rest:
                rp = p.add_run()
                rp.text           = rest
                rp.font.name      = FONT
                rp.font.size      = Pt(pt)
                rp.font.bold      = False
                rp.font.color.rgb = C["text_dark"]
        else:
            r1 = p.add_run()
            r1.text           = item["text"]
            r1.font.name      = FONT
            r1.font.size      = Pt(pt)
            r1.font.bold      = False
            r1.font.color.rgb = C["text_dark"]

# ── ficha tabla ────────────────────────────────────────────────────────────

def ficha_table(slide, ficha: dict, x, y, w, h) -> None:
    rows = list(ficha.items())
    n    = len(rows)

    # fondo
    rect(slide, x, y, w, h, C["primary"])

    # header label
    tb = textbox(slide, x + 10, y + 14, w - 20, 18)
    para_first(tb.text_frame, "FICHA DEL MÓDULO", 8, bold=True,
               color=RGBColor(180, 215, 225))

    # tabla de datos
    TABLE_TOP = y + 44
    TABLE_H   = h - 54
    ROW_H     = TABLE_H // n

    t = slide.shapes.add_table(n, 2, px(x), px(TABLE_TOP), px(w), px(TABLE_H)).table
    t.columns[0].width = px(w * 0.55)
    t.columns[1].width = px(w * 0.45)

    for i, (label, val) in enumerate(rows):
        row        = t.rows[i]
        row.height = px(ROW_H)

        lc = row.cells[0]
        vc = row.cells[1]

        cell_solid(lc, C["primary_dark"])
        cell_solid(vc, C["primary"])
        cell_no_border(lc)
        cell_no_border(vc)

        cell_write(lc, label, 7.5, bold=True, color=RGBColor(180,215,225), margin_top=8)
        cell_write(vc, val,   15,  bold=True, color=C["white"],             margin_top=6)

# ── parseo HTML ────────────────────────────────────────────────────────────

def parse_html(path: Path) -> list[dict]:
    soup   = BeautifulSoup(path.read_text("utf-8"), "lxml")
    result = []
    for sec in soup.find_all("section", class_="slide-container"):
        result.append(_classify(sec))
    return result

def _txt(el) -> str:
    return el.get_text(" ", strip=True) if el else ""

def _classify(sec: Tag) -> dict:
    h2  = sec.find("h2", class_="slide-title")
    num = sec.find(class_="slide-number")
    n   = int(num.get_text(strip=True)) if num else 0

    if not h2:
        # title-layout
        h1    = sec.find("h1")
        title = _txt(h1) if h1 else ""
        return {"type":"cover","title":title,"n":n}

    title = _txt(h2)

    if re.match(r"^Módulo \d+", title):
        return _parse_module(sec, title, n)

    if "Competencias INAP" in title:
        return _parse_benchmark(sec, title, n)

    if "quién va" in title or "Perfiles" in title:
        return _parse_profiles(sec, title, n)

    if "Metodología" in title:
        return _parse_metodologia(sec, title, n)

    if "Propuesta de valor" in title:
        return _parse_propuesta(sec, title, n)

    if "visión" in title.lower():
        return _parse_vision(sec, title, n)

    if "Por qué" in title or "por qué" in title.lower() or "reto" in title.lower():
        return _parse_porque(sec, title, n)

    # genérico
    return {"type":"generic","title":title,"n":n}


def _parse_module(sec: Tag, title: str, n: int) -> dict:
    h3       = sec.find("h3")
    subtitle = _txt(h3)
    items    = []
    for li in sec.select(".topics-list li"):
        strongs = [s.get_text(strip=True) for s in li.find_all("strong")]
        items.append({"text": _txt(li), "bold_prefix": strongs[0] if strongs else None})
    return {"type":"module","title":title,"subtitle":subtitle,"items":items,"n":n,
            "ficha":{"HORAS DE FORMACIÓN":"—","SESIONES ONLINE":"—",
                     "SESIONES PRESENCIALES":"—","CLAUSTRO":"—"}}


def _parse_benchmark(sec: Tag, title: str, n: int) -> dict:
    rows = []
    for row in sec.select(".benchmark-row:not(.header)"):
        cols = [c.get_text(strip=True) for c in row.find_all("div", recursive=False)]
        if cols:
            rows.append(cols)
    return {"type":"benchmark","title":title,"rows":rows,"n":n}


def _parse_profiles(sec: Tag, title: str, n: int) -> dict:
    cards = []
    for card in sec.select(".profile-card"):
        h4  = card.find("h4")
        p   = card.find("p")
        cards.append({"h": _txt(h4), "p": _txt(p)})
    return {"type":"profiles","title":title,"cards":cards,"n":n}


def _parse_metodologia(sec: Tag, title: str, n: int) -> dict:
    cards = []
    for div in sec.select(".content-area > div > div"):
        h4 = div.find("h4")
        p  = div.find("p")
        em = div.find("div", style=lambda s: s and "font-size: 32px" in (s or ""))
        cards.append({"icon": _txt(em), "h": _txt(h4), "p": _txt(p)})
    return {"type":"metodologia","title":title,"cards":cards,"n":n}


def _parse_propuesta(sec: Tag, title: str, n: int) -> dict:
    cards = []
    for card in sec.select(".value-card"):
        num = card.find(class_="value-num")
        h4  = card.find("h4")
        p   = card.find("p")
        cards.append({"num": _txt(num), "h": _txt(h4), "p": _txt(p)})
    footer = sec.find(style=lambda s: s and "text-align: center" in (s or ""))
    return {"type":"propuesta","title":title,"cards":cards,"n":n,
            "footer": _txt(footer) if footer else ""}


def _parse_vision(sec: Tag, title: str, n: int) -> dict:
    pillars = []
    for div in sec.select(".content-area > div > div[style]"):
        h3 = div.find("h3")
        p  = div.find("p")
        if h3:
            pillars.append({"h": _txt(h3), "p": _txt(p)})
    return {"type":"vision","title":title,"pillars":pillars,"n":n}


def _parse_porque(sec: Tag, title: str, n: int) -> dict:
    intro = sec.find("p", style=lambda s: s and "font-size: 18px" in (s or ""))
    stats = []
    for card in sec.select(".stat-card"):
        num = card.find(class_="stat-num")
        lbl = card.find(class_="stat-label")
        stats.append({"num": _txt(num), "label": _txt(lbl)})
    h3 = sec.find("h3")
    return {"type":"porque","title":title,"h3": _txt(h3) if h3 else "",
            "intro": _txt(intro),"stats":stats,"n":n}


# ── renderizadores por tipo ────────────────────────────────────────────────

def render_cover(slide, data: dict) -> None:
    rect(slide, 0, 0, W, H, C["bg_dark"])
    # franja de color en la parte inferior
    rect(slide, 0, H - 8, W, 8, C["primary"])
    rect(slide, 0, H//2 - 3, W, 6, C["primary"])

    # título central
    title = data["title"]
    # buscar si hay span (texto en color)
    # El HTML tiene <h1>... <span>digital</span></h1>
    # Aproximamos con dos runs
    parts = re.split(r"\s*\n\s*", title.strip())
    tb = textbox(slide, 80, H//2 - 110, W - 160, 200)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)
    add_run(p, title, 52, bold=True, color=C["white"])

    # subtítulo
    subtitle = data.get("subtitle", "")
    if not subtitle:
        # obtener del HTML original si hay
        subtitle = "Certificación oficial · esPublico"  # portada default
    if data["n"] == 1:
        tb2 = textbox(slide, 80, H//2 + 100, W - 160, 40)
        para_first(tb2.text_frame, subtitle, 14, color=C["accent"])
    slide_number_label(slide, data["n"])


def render_cover_alt(slide, data: dict) -> None:
    """Cierre."""
    rect(slide, 0, 0, W, H, C["bg_dark"])
    rect(slide, 0, H - 8, W, 8, C["primary"])
    tb = textbox(slide, 80, H//2 - 80, W - 160, 160)
    tf = tb.text_frame
    tf.word_wrap = True
    para_first(tf, data["title"], 44, bold=True, color=C["white"])
    slide_number_label(slide, data["n"])


def render_module(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"], y=16, w_offset=400)
    divider(slide, y=62, w=W - 450)

    # subtítulo
    if data.get("subtitle"):
        tb = textbox(slide, 40, 72, W - 440, 32)
        para_first(tb.text_frame, data["subtitle"], 13, bold=True, color=C["text_dark"])

    # bullets
    bullet_list(slide, data["items"], x=40, y=112, w=W - 440, h=H - 130, pt=11)

    # ficha tabla
    FICHA_X, FICHA_W = W - 390, 370
    ficha_table(slide, data["ficha"], FICHA_X, 10, FICHA_W, H - 20)

    slide_number_label(slide, data["n"])


def render_porque(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"])
    divider(slide)

    # h3
    if data.get("h3"):
        tb = textbox(slide, 40, 76, W - 80, 32)
        para_first(tb.text_frame, data["h3"], 14, bold=True, color=C["text_dark"])

    # intro
    if data.get("intro"):
        tb2 = textbox(slide, 40, 110, W - 80, 70)
        tf  = tb2.text_frame
        tf.word_wrap = True
        para_first(tf, data["intro"], 12, color=C["text_body"])

    # stats cards (3 columnas)
    stats   = data["stats"]
    CARD_W  = (W - 80 - 32) // 3
    CARD_H  = 160
    CARD_Y  = 192
    for i, s in enumerate(stats):
        cx = 40 + i * (CARD_W + 16)
        rect(slide, cx, CARD_Y, CARD_W, CARD_H, C["bg_light"], rounded=True)
        # borde superior
        rect(slide, cx, CARD_Y, CARD_W, 4, C["primary"])

        tb = textbox(slide, cx + 16, CARD_Y + 14, CARD_W - 32, 50)
        para_first(tb.text_frame, s["num"], 28, bold=True, color=C["primary"])

        tb2 = textbox(slide, cx + 16, CARD_Y + 66, CARD_W - 32, 80)
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        para_first(tf2, s["label"], 11, color=C["text_muted"])

    slide_number_label(slide, data["n"])


def render_vision(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"])
    divider(slide)

    pillars = data.get("pillars", [])
    if len(pillars) >= 2:
        COL_W = (W - 80 - 48 - 60) // 2   # espacio para el "+"
        COL_Y = 80
        COL_H = H - COL_Y - 20

        # pilar izquierdo — oscuro
        rect(slide, 40, COL_Y, COL_W, COL_H, C["bg_dark"], rounded=True)
        tb1 = textbox(slide, 72, COL_Y + 28, COL_W - 64, 36)
        para_first(tb1.text_frame, pillars[0]["h"], 17, bold=True, color=C["white"])
        tb2 = textbox(slide, 72, COL_Y + 72, COL_W - 64, COL_H - 100)
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        para_first(tf2, pillars[0]["p"], 12, color=RGBColor(230,240,240))

        # símbolo "+"
        MID_X = 40 + COL_W + 4
        tb_plus = textbox(slide, MID_X, H//2 - 30, 52, 60)
        para_first(tb_plus.text_frame, "+", 40, bold=False, color=C["primary"],
                   align=PP_ALIGN.CENTER)

        # pilar derecho — primario
        COL2_X = 40 + COL_W + 60
        rect(slide, COL2_X, COL_Y, COL_W, COL_H, C["primary"], rounded=True)
        tb3 = textbox(slide, COL2_X + 32, COL_Y + 28, COL_W - 64, 36)
        para_first(tb3.text_frame, pillars[1]["h"], 17, bold=True, color=C["white"])
        tb4 = textbox(slide, COL2_X + 32, COL_Y + 72, COL_W - 64, COL_H - 100)
        tf4 = tb4.text_frame
        tf4.word_wrap = True
        para_first(tf4, pillars[1]["p"], 12, color=RGBColor(230,244,244))

    slide_number_label(slide, data["n"])


def render_propuesta(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"])
    divider(slide)

    cards  = data["cards"]
    n_cols = 4
    n_rows = (len(cards) + n_cols - 1) // n_cols
    CARD_W = (W - 80 - (n_cols - 1) * 16) // n_cols
    AREA_H = H - 88 - 36 - (12 if data.get("footer") else 0)
    CARD_H = (AREA_H - (n_rows - 1) * 12) // n_rows
    CARD_Y0 = 80

    for i, card in enumerate(cards):
        row = i // n_cols
        col = i % n_cols
        cx  = 40 + col * (CARD_W + 16)
        cy  = CARD_Y0 + row * (CARD_H + 12)

        rect(slide, cx, cy, CARD_W, CARD_H, C["bg_dark"], rounded=True)

        # número
        tb_n = textbox(slide, cx + 16, cy + 12, CARD_W - 32, 32)
        para_first(tb_n.text_frame, card["num"], 20, bold=True, color=C["accent"])

        # titulo
        tb_h = textbox(slide, cx + 16, cy + 46, CARD_W - 32, 30)
        para_first(tb_h.text_frame, card["h"], 11.5, bold=True, color=C["white"])

        # descripcion
        tb_p = textbox(slide, cx + 16, cy + 78, CARD_W - 32, CARD_H - 90)
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        para_first(tf_p, card["p"], 9.5, color=RGBColor(200,220,225))

    # footer
    if data.get("footer"):
        tb_f = textbox(slide, 40, H - 38, W - 80, 26)
        tf_f = tb_f.text_frame
        tf_f.word_wrap = True
        rect(slide, 40, H - 44, W - 80, 32, C["bg_light"], rounded=True)
        para_first(tb_f.text_frame, data["footer"], 10, color=C["text_dark"],
                   align=PP_ALIGN.CENTER)

    slide_number_label(slide, data["n"])


def render_metodologia(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"])
    divider(slide)

    cards  = data["cards"]
    n_cols = 3
    n_rows = (len(cards) + n_cols - 1) // n_cols
    CARD_W = (W - 80 - (n_cols - 1) * 16) // n_cols
    AREA_H = H - 80 - 20
    CARD_H = (AREA_H - (n_rows - 1) * 12) // n_rows
    CARD_Y0 = 78

    for i, card in enumerate(cards):
        row = i // n_cols
        col = i % n_cols
        cx  = 40 + col * (CARD_W + 16)
        cy  = CARD_Y0 + row * (CARD_H + 12)

        rect(slide, cx, cy, CARD_W, CARD_H, C["white"], rounded=True)
        # borde
        s = slide.shapes.add_shape(1, px(cx), px(cy), px(CARD_W), px(CARD_H))
        s.fill.background()
        s.line.color.rgb = C["border"]
        s.line.width     = Pt(0.75)

        # icono
        if card.get("icon"):
            tb_i = textbox(slide, cx + 16, cy + 12, 40, 36)
            para_first(tb_i.text_frame, card["icon"], 22)

        # título
        tb_h = textbox(slide, cx + 16, cy + 50, CARD_W - 32, 28)
        para_first(tb_h.text_frame, card["h"], 11, bold=True, color=C["text_dark"])

        # descripción
        tb_p = textbox(slide, cx + 16, cy + 80, CARD_W - 32, CARD_H - 96)
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        para_first(tf_p, card["p"], 10, color=C["text_muted"])

    slide_number_label(slide, data["n"])


def render_benchmark(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"])
    divider(slide)

    rows     = data["rows"]
    n_rows   = len(rows) + 1   # +1 header
    TABLE_Y  = 74
    TABLE_H  = H - TABLE_Y - 24
    TABLE_W  = W - 80
    ROW_H    = TABLE_H // n_rows

    # columnas: área 1fr, competencia 1.5fr, módulos 1fr → proporciones
    COL_W = [int(TABLE_W * 0.22), int(TABLE_W * 0.55), int(TABLE_W * 0.23)]

    t = slide.shapes.add_table(
        n_rows, 3, px(40), px(TABLE_Y), px(TABLE_W), px(TABLE_H)
    ).table
    t.columns[0].width = px(COL_W[0])
    t.columns[1].width = px(COL_W[1])
    t.columns[2].width = px(COL_W[2])

    # header
    for ci, label in enumerate(["Área", "Competencia", "Módulos"]):
        c = t.rows[0].cells[ci]
        t.rows[0].height = px(ROW_H)
        cell_solid(c, C["primary"])
        cell_no_border(c)
        cell_write(c, label, 10, bold=True, color=C["white"], margin_top=8)

    # filas
    for ri, row_data in enumerate(rows):
        row        = t.rows[ri + 1]
        row.height = px(ROW_H)
        bg_color   = C["white"] if ri % 2 == 0 else C["bg_light"]
        for ci, val in enumerate(row_data[:3]):
            c = row.cells[ci]
            cell_solid(c, bg_color)
            cell_no_border(c)
            bold = ci == 0
            cell_write(c, val, 9.5, bold=bold,
                       color=C["primary"] if ci == 0 else C["text_dark"],
                       margin_top=6)

    slide_number_label(slide, data["n"])


def render_profiles(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"])
    divider(slide)

    cards  = data["cards"]
    n_cols = 2
    n_rows = (len(cards) + n_cols - 1) // n_cols
    CARD_W = (W - 80 - 16) // n_cols
    AREA_H = H - 80 - 20
    CARD_H = (AREA_H - (n_rows - 1) * 12) // n_rows
    CARD_Y0 = 78

    for i, card in enumerate(cards):
        row = i // n_cols
        col = i % n_cols
        cx  = 40 + col * (CARD_W + 16)
        cy  = CARD_Y0 + row * (CARD_H + 12)

        s = slide.shapes.add_shape(1, px(cx), px(cy), px(CARD_W), px(CARD_H))
        s.fill.background()
        s.line.color.rgb = C["border"]
        s.line.width     = Pt(0.75)

        tb_h = textbox(slide, cx + 16, cy + 14, CARD_W - 32, 28)
        para_first(tb_h.text_frame, card["h"], 12, bold=True, color=C["primary"])

        tb_p = textbox(slide, cx + 16, cy + 44, CARD_W - 32, CARD_H - 58)
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        para_first(tf_p, card["p"], 10.5, color=C["text_muted"])

    slide_number_label(slide, data["n"])


def render_generic(slide, data: dict) -> None:
    slide_bg(slide)
    slide_top_bar(slide)
    slide_title_text(slide, data["title"])
    divider(slide)
    slide_number_label(slide, data["n"])


# ── dispatcher ─────────────────────────────────────────────────────────────

RENDERERS = {
    "cover":       render_cover,
    "module":      render_module,
    "porque":      render_porque,
    "vision":      render_vision,
    "propuesta":   render_propuesta,
    "metodologia": render_metodologia,
    "benchmark":   render_benchmark,
    "profiles":    render_profiles,
    "generic":     render_generic,
}

def render(slide, data: dict) -> None:
    t = data["type"]
    if t == "cover" and data["n"] > 1:
        render_cover_alt(slide, data)
    else:
        RENDERERS.get(t, render_generic)(slide, data)


# ── construir PPTX ─────────────────────────────────────────────────────────

def build_pptx(slides_data: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width  = Emu(W_EMU)
    prs.slide_height = Emu(H_EMU)
    blank = prs.slide_layouts[6]

    for i, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        render(slide, data)
        print(f"[pptx] slide {i+1:02d} ({data['type']})", file=sys.stderr)

    prs.save(OUT_PPTX)
    print(f"[done] {OUT_PPTX.name} ({OUT_PPTX.stat().st_size // 1024} KB)")


# ── PDF ────────────────────────────────────────────────────────────────────

def export_pdf() -> None:
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx  = browser.new_context(viewport={"width": W, "height": H})
        page = ctx.new_page()
        page.goto(INDEX.as_uri(), wait_until="networkidle")
        page.emulate_media(media="print")
        page.wait_for_timeout(800)
        page.pdf(path=str(OUT_PDF), width=f"{W}px", height=f"{H}px",
                 print_background=True,
                 margin={"top":"0","right":"0","bottom":"0","left":"0"})
        browser.close()
    print(f"[done] {OUT_PDF.name} ({OUT_PDF.stat().st_size // 1024} KB)")


# ── main ───────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("[parse] leyendo HTML …", file=sys.stderr)
    slides_data = parse_html(INDEX)
    for d in slides_data:
        print(f"  slide {d['n']:02d}  type={d['type']:12s}  {d.get('title','')[:50]}", file=sys.stderr)

    print("[pptx] construyendo …", file=sys.stderr)
    build_pptx(slides_data)

    print("[pdf] generando …", file=sys.stderr)
    export_pdf()
