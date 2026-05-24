"""
build_pptx_editable.py — PPTX editable + PDF desde index.html.

Estrategia:
  - Todas las slides: screenshot full-bleed como fondo (preserva diseño exacto)
  - Slides de módulo: tabla editable superpuesta sobre el área de la ficha
    usando coordenadas exactas obtenidas con Playwright bounding_box()
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from bs4 import BeautifulSoup, Tag
from playwright.sync_api import sync_playwright, Page
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ── rutas ──────────────────────────────────────────────────────────────────
ROOT      = Path(__file__).resolve().parent.parent
INDEX     = ROOT / "index.html"
SHOTS_DIR = ROOT / ".tools" / "shots"
OUT_PPTX  = ROOT / "certificacion-inap.pptx"
OUT_PDF   = ROOT / "certificacion-inap.pdf"

# ── dimensiones ────────────────────────────────────────────────────────────
W_PX, H_PX = 1280, 720
EMU_PX      = 9525
W_EMU       = W_PX * EMU_PX
H_EMU       = H_PX * EMU_PX

# ── paleta gestiona ────────────────────────────────────────────────────────
PRIMARY      = RGBColor(0,   109, 133)   # #006d85
PRIMARY_DARK = RGBColor(0,    77,  94)   # #004d5e
WHITE        = RGBColor(255, 255, 255)
TEXT_MUTED   = RGBColor(200, 225, 230)
FONT         = "Roboto"


# ── helpers EMU ────────────────────────────────────────────────────────────

def px(n: float) -> Emu:
    return Emu(int(n * EMU_PX))


# ── parseo HTML: extraer datos de ficha de cada módulo ────────────────────

def parse_module_fichas(html_path: Path) -> dict[int, dict]:
    """Devuelve {slide_idx: ficha_data} para slides de módulo."""
    soup   = BeautifulSoup(html_path.read_text(encoding="utf-8"), "lxml")
    result: dict[int, dict] = {}
    for i, sec in enumerate(soup.find_all("section", class_="slide-container")):
        title_el = sec.find("h2", class_="slide-title")
        if not title_el:
            continue
        title = title_el.get_text(strip=True)
        if re.match(r"^Módulo \d+", title):
            result[i] = {
                "HORAS DE FORMACIÓN":    "—",
                "SESIONES ONLINE":       "—",
                "SESIONES PRESENCIALES": "—",
                "CLAUSTRO":              "—",
            }
    return result


# ── capturas con Playwright ────────────────────────────────────────────────

def capture_all(module_indices: set[int]) -> tuple[dict[int, Path], dict[int, dict]]:
    """
    Captura screenshot de cada slide.
    Para slides de módulo, también obtiene bounding_box de la ficha.
    Devuelve (screenshots, ficha_boxes).
    ficha_boxes: {slide_idx: {x, y, w, h}} en píxeles de pantalla.
    """
    SHOTS_DIR.mkdir(parents=True, exist_ok=True)
    screenshots: dict[int, Path] = {}
    ficha_boxes: dict[int, dict] = {}

    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx     = browser.new_context(
            viewport={"width": W_PX, "height": H_PX},
            device_scale_factor=2,          # 2× para calidad
        )
        page = ctx.new_page()
        page.goto(INDEX.as_uri(), wait_until="networkidle")
        page.emulate_media(media="print")
        page.wait_for_timeout(500)

        containers = page.locator(".slide-container")
        total      = containers.count()
        print(f"[info] {total} slides encontradas", file=sys.stderr)

        for i in range(total):
            target = SHOTS_DIR / f"slide_{i+1:02d}.png"
            containers.nth(i).screenshot(path=str(target))
            screenshots[i] = target
            print(f"[shot] slide {i+1:02d}", file=sys.stderr)

            if i in module_indices:
                # bounding_box de la columna ficha (último div hijo directo de content-area)
                # buscamos el div con background color-primary dentro del slide
                ficha_loc = containers.nth(i).locator(
                    "div[style*='background: var(--color-primary)']"
                ).first
                try:
                    bb = ficha_loc.bounding_box()
                    if bb:
                        # bb viene en coords relativas a la ventana (viewport)
                        # como capturamos el elemento, necesitamos relativas al slide
                        slide_bb = containers.nth(i).bounding_box()
                        ficha_boxes[i] = {
                            "x": bb["x"] - (slide_bb["x"] if slide_bb else 0),
                            "y": bb["y"] - (slide_bb["y"] if slide_bb else 0),
                            "w": bb["width"],
                            "h": bb["height"],
                        }
                        print(f"[bbox] slide {i+1:02d} ficha: {ficha_boxes[i]}", file=sys.stderr)
                except Exception as e:
                    print(f"[warn] slide {i+1:02d} bbox error: {e}", file=sys.stderr)

        browser.close()
    return screenshots, ficha_boxes


# ── tabla editable de ficha ────────────────────────────────────────────────

def _cell_fill(cell, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # eliminar fills previos si existen
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    sf     = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb   = etree.SubElement(sf,   qn("a:srgbClr"))
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def _cell_text(cell, text: str, pt: float, bold=False, color=WHITE) -> None:
    tf = cell.text_frame
    tf.clear()
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.color.rgb = color


def add_ficha_table(slide, ficha: dict, box: dict) -> None:
    """Superpone tabla editable encima del área de la ficha."""
    x, y, w, h = box["x"], box["y"], box["w"], box["h"]

    # Ajuste device_scale_factor=2: coordenadas ya están en px físicos,
    # pero bounding_box devuelve CSS px (1×), así que no dividimos.
    rows_data = list(ficha.items())
    n_rows    = len(rows_data)

    # padding superior para no tapar el header "FICHA DEL MÓDULO"
    # (aprox 48px del header + margen)
    HEADER_H = 52
    table_y  = y + HEADER_H
    table_h  = h - HEADER_H - 4

    tbl_shape = slide.shapes.add_table(
        n_rows, 2,
        px(x), px(table_y),
        px(w), px(table_h),
    )
    tbl = tbl_shape.table

    col_w_label = int(w * 0.58)
    col_w_value = w - col_w_label
    tbl.columns[0].width = px(col_w_label)
    tbl.columns[1].width = px(col_w_value)

    row_h = table_h // n_rows

    for ri, (label, value) in enumerate(rows_data):
        row         = tbl.rows[ri]
        row.height  = px(row_h)
        label_cell  = row.cells[0]
        value_cell  = row.cells[1]

        _cell_fill(label_cell, PRIMARY_DARK)
        _cell_fill(value_cell, PRIMARY)

        _cell_text(label_cell, label, 7.5, bold=True, color=TEXT_MUTED)
        _cell_text(value_cell, value, 15,  bold=True, color=WHITE)

        # sin bordes
        for cell in (label_cell, value_cell):
            _remove_cell_borders(cell)


def _remove_cell_borders(cell) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(border))
        if existing is None:
            existing = etree.SubElement(tcPr, qn(border))
        # limpiar hijos y poner noFill
        for child in list(existing):
            existing.remove(child)
        etree.SubElement(existing, qn("a:noFill"))


# ── construir PPTX ─────────────────────────────────────────────────────────

def build_pptx(
    n_slides:  int,
    fichas:    dict[int, dict],
    shots:     dict[int, Path],
    boxes:     dict[int, dict],
) -> None:
    prs = Presentation()
    prs.slide_width  = Emu(W_EMU)
    prs.slide_height = Emu(H_EMU)
    blank = prs.slide_layouts[6]

    for i in range(n_slides):
        slide = prs.slides.add_slide(blank)

        # 1. fondo: screenshot completo (preserva diseño)
        img = shots.get(i)
        if img and img.exists():
            slide.shapes.add_picture(str(img), 0, 0, Emu(W_EMU), Emu(H_EMU))

        # 2. si es módulo: tabla editable sobre la ficha
        if i in fichas and i in boxes:
            add_ficha_table(slide, fichas[i], boxes[i])
            print(f"[pptx] slide {i+1:02d} módulo (img+tabla editable)", file=sys.stderr)
        else:
            print(f"[pptx] slide {i+1:02d} imagen", file=sys.stderr)

    prs.save(OUT_PPTX)
    print(f"[done] {OUT_PPTX.name} ({OUT_PPTX.stat().st_size // 1024} KB)")


# ── PDF ────────────────────────────────────────────────────────────────────

def export_pdf() -> None:
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx  = browser.new_context(viewport={"width": W_PX, "height": H_PX})
        page = ctx.new_page()
        page.goto(INDEX.as_uri(), wait_until="networkidle")
        page.emulate_media(media="print")
        page.wait_for_timeout(800)
        page.pdf(
            path=str(OUT_PDF),
            width=f"{W_PX}px",
            height=f"{H_PX}px",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )
        browser.close()
    print(f"[done] {OUT_PDF.name} ({OUT_PDF.stat().st_size // 1024} KB)")


# ── main ───────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("[parse] identificando módulos …", file=sys.stderr)
    fichas_data = parse_module_fichas(INDEX)
    print(f"[parse] módulos en slides: {sorted(fichas_data.keys())}", file=sys.stderr)

    print("[browser] capturas + bounding boxes …", file=sys.stderr)
    screenshots, ficha_boxes = capture_all(set(fichas_data.keys()))

    print("[pptx] construyendo …", file=sys.stderr)
    build_pptx(len(screenshots), fichas_data, screenshots, ficha_boxes)

    print("[pdf] generando …", file=sys.stderr)
    export_pdf()
