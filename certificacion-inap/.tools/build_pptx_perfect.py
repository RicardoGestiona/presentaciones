"""
build_pptx_perfect.py — PPTX 100 % nativo, fiel al render HTML.

Walker DOM extrae cada elemento visible como átomo (shape o text) con
posición, estilo y formato calculado por el navegador. python-pptx los
reconstruye uno a uno → cada elemento es seleccionable.
"""
from __future__ import annotations

import sys
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

ROOT     = Path(__file__).resolve().parent.parent
INDEX    = ROOT / "index.html"
OUT_PPTX = ROOT / "certificacion-inap.pptx"

W_CSS, H_CSS = 1280, 720
EMU_PX       = 9525
W_EMU        = W_CSS * EMU_PX
H_EMU        = H_CSS * EMU_PX
FONT         = "Roboto"

def px(n: float) -> Emu:
    return Emu(int(round(n * EMU_PX)))


# ─── JS: extracción de átomos ─────────────────────────────────────────────
JS_EXTRACT = r"""
(slideIndex) => {
    const slide = document.querySelectorAll('.slide-container')[slideIndex];
    if (!slide) return [];
    const sr = slide.getBoundingClientRect();
    const atoms = [];

    const SKIP_TAGS = new Set(['SCRIPT','STYLE','LINK','META','BR','I','EM']);

    function rgb(s) {
        const m = s && s.match(/(\d+(\.\d+)?)/g);
        if (!m || m.length < 3) return null;
        return [parseInt(m[0]), parseInt(m[1]), parseInt(m[2])];
    }
    function alphaOf(s) {
        if (!s) return 1;
        const m = s.match(/rgba?\(([^)]+)\)/);
        if (!m) return 1;
        const parts = m[1].split(',').map(p => parseFloat(p.trim()));
        return parts.length === 4 ? parts[3] : 1;
    }

    function pseudo(el, which) {
        try {
            const st = window.getComputedStyle(el, '::' + which);
            let c = st.content || '';
            if (c === 'none' || c === 'normal' || c === '""' || c === "''") return null;
            c = c.replace(/^["']|["']$/g, '').replace(/\\([0-9a-fA-F]+)\s?/g, (m, hex) => String.fromCodePoint(parseInt(hex,16)));
            return {
                text: c,
                color: rgb(st.color),
                fontSize: parseFloat(st.fontSize),
                fontWeight: parseInt(st.fontWeight) || 400
            };
        } catch(e) { return null; }
    }

    function runs(el) {
        const out = [];
        const baseSt = window.getComputedStyle(el);
        function walk(node, ctx) {
            if (node.nodeType === 3) {
                const t = node.textContent;
                if (t && t.length) out.push({ text: t, bold: ctx.bold, italic: ctx.italic, color: ctx.color });
            } else if (node.nodeType === 1) {
                const tag = node.tagName.toLowerCase();
                if (tag === 'br') { out.push({ text: '\n', bold: ctx.bold, italic: ctx.italic, color: ctx.color }); return; }
                const st = window.getComputedStyle(node);
                const newCtx = {
                    bold: (parseInt(st.fontWeight) || 400) >= 600 || tag === 'strong' || tag === 'b' || ctx.bold,
                    italic: st.fontStyle === 'italic' || tag === 'em' || tag === 'i' || ctx.italic,
                    color: rgb(st.color) || ctx.color
                };
                for (const c of node.childNodes) walk(c, newCtx);
            }
        }
        walk(el, {
            bold: (parseInt(baseSt.fontWeight) || 400) >= 600,
            italic: baseSt.fontStyle === 'italic',
            color: rgb(baseSt.color)
        });
        return out;
    }

    function walk(el) {
        if (SKIP_TAGS.has(el.tagName)) return;

        const st = window.getComputedStyle(el);
        if (st.display === 'none' || st.visibility === 'hidden') return;

        const r = el.getBoundingClientRect();
        const x = r.left - sr.left, y = r.top - sr.top;
        const w = r.width, h = r.height;
        if (w < 0.5 || h < 0.5) {
            for (const c of el.children) walk(c);
            return;
        }

        // SHAPE atom: emite si hay fondo, borde o radio
        const bg = rgb(st.backgroundColor);
        const bgA = alphaOf(st.backgroundColor);
        const hasBg = bg && bgA > 0.01;

        const bt = parseFloat(st.borderTopWidth) || 0;
        const br = parseFloat(st.borderRightWidth) || 0;
        const bb = parseFloat(st.borderBottomWidth) || 0;
        const bl = parseFloat(st.borderLeftWidth) || 0;
        const hasBorder = bt + br + bb + bl > 0.5;

        const rad = parseFloat(st.borderTopLeftRadius) || 0;

        if (hasBg || hasBorder) {
            atoms.push({
                kind: 'shape',
                x, y, w, h,
                bg:      hasBg ? bg : null,
                bgAlpha: bgA,
                bt, btC: rgb(st.borderTopColor),
                br, brC: rgb(st.borderRightColor),
                bb, bbC: rgb(st.borderBottomColor),
                bl, blC: rgb(st.borderLeftColor),
                radius:  rad
            });
        }

        // detectar leaf de texto (sin hijos bloque)
        const children = Array.from(el.children).filter(c => window.getComputedStyle(c).display !== 'none');
        const blockChildren = children.filter(c => {
            const d = window.getComputedStyle(c).display;
            return d !== 'inline' && d !== 'inline-block' && d !== 'inline-flex' && d !== 'contents';
        });

        const hasTextOwn = el.innerText && el.innerText.trim().length > 0;

        if (blockChildren.length === 0 && hasTextOwn) {
            const pre = pseudo(el, 'before');
            const post = pseudo(el, 'after');
            const r0 = runs(el);
            // limpiar columns CSS — no se soporta en PPTX, usa solo 1 col
            let lh = st.lineHeight;
            let lhNum = parseFloat(lh);
            if (isNaN(lhNum) || lh === 'normal') lhNum = parseFloat(st.fontSize) * 1.2;

            atoms.push({
                kind: 'text',
                x, y, w, h,
                runs: r0,
                before: pre, after: post,
                fontSize: parseFloat(st.fontSize),
                fontWeight: parseInt(st.fontWeight) || 400,
                color: rgb(st.color),
                align: st.textAlign,
                italic: st.fontStyle === 'italic',
                lineHeight: lhNum,
                paddingLeft: parseFloat(st.paddingLeft) || 0,
                paddingTop: parseFloat(st.paddingTop) || 0,
                paddingRight: parseFloat(st.paddingRight) || 0,
                paddingBottom: parseFloat(st.paddingBottom) || 0,
                letterSpacing: parseFloat(st.letterSpacing) || 0,
                textTransform: st.textTransform,
                uppercase: st.textTransform === 'uppercase'
            });
        } else {
            for (const c of children) walk(c);
        }
    }
    walk(slide);
    return atoms;
}
"""


# ─── PPTX helpers ─────────────────────────────────────────────────────────

def _make_transparent(shape) -> None:
    """noFill real (en XML)."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None: return
    for tag in ("a:solidFill","a:gradFill","a:pattFill","a:blipFill","a:grpFill","a:noFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    etree.SubElement(spPr, qn("a:noFill"))


def _set_line(shape, color, width_pt: float) -> None:
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for ch in list(ln): ln.remove(ch)
    ln.set("w", str(int(width_pt * 12700)))   # EMU per pt
    sf = etree.SubElement(ln, qn("a:solidFill"))
    srgb = etree.SubElement(sf, qn("a:srgbClr"))
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def _no_line(shape) -> None:
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for ch in list(ln): ln.remove(ch)
    etree.SubElement(ln, qn("a:noFill"))


def _align(a: str) -> PP_ALIGN:
    return {
        "left":    PP_ALIGN.LEFT,
        "center":  PP_ALIGN.CENTER,
        "right":   PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
        "start":   PP_ALIGN.LEFT,
        "end":     PP_ALIGN.RIGHT,
    }.get(a, PP_ALIGN.LEFT)


def add_shape_atom(slide, a: dict) -> None:
    """Crea rectángulo (con radius opcional) + opcional borde."""
    w_min = min(a["w"], a["h"])
    rounded = a["radius"] > 0 and w_min > 0

    if rounded:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    px(a["x"]), px(a["y"]), px(a["w"]), px(a["h"]))
        # adjustment[0] es la proporción del radio sobre la mitad del lado más corto
        s.adjustments[0] = min(a["radius"] / w_min, 0.5)
    else:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    px(a["x"]), px(a["y"]), px(a["w"]), px(a["h"]))

    # Fill
    if a["bg"]:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(*a["bg"])
        # alpha
        if 0 < a["bgAlpha"] < 1:
            srgb = s.fill.fore_color._xFill.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", str(int(a["bgAlpha"] * 100000)))
    else:
        _make_transparent(s)

    # Border (CSS bordes por lado; PPTX solo soporta un borde uniforme).
    # Estrategia: si los 4 bordes son iguales → línea uniforme.
    # Si solo un lado tiene borde → emitir línea adicional como rectángulo fino.
    sides = [(a["bt"], a["btC"]), (a["br"], a["brC"]),
             (a["bb"], a["bbC"]), (a["bl"], a["blC"])]
    non_zero = [(w_,c_) for w_,c_ in sides if w_ > 0.5 and c_]

    if len(non_zero) == 4 and len(set((round(w_,1), tuple(c_)) for w_,c_ in non_zero)) == 1:
        # 4 bordes iguales
        w_, c_ = non_zero[0]
        _set_line(s, c_, w_ * 0.75)
    else:
        _no_line(s)
        # bordes parciales como líneas separadas (border-top, border-left, etc.)
        if a["bt"] > 0.5 and a["btC"]:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          px(a["x"]), px(a["y"]),
                                          px(a["w"]), px(a["bt"]))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*a["btC"])
            _no_line(line)
        if a["bb"] > 0.5 and a["bbC"]:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          px(a["x"]), px(a["y"] + a["h"] - a["bb"]),
                                          px(a["w"]), px(a["bb"]))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*a["bbC"])
            _no_line(line)
        if a["bl"] > 0.5 and a["blC"]:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          px(a["x"]), px(a["y"]),
                                          px(a["bl"]), px(a["h"]))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*a["blC"])
            _no_line(line)
        if a["br"] > 0.5 and a["brC"]:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          px(a["x"] + a["w"] - a["br"]), px(a["y"]),
                                          px(a["br"]), px(a["h"]))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*a["brC"])
            _no_line(line)


def add_text_atom(slide, a: dict) -> None:
    """Textbox transparente con runs formateados."""
    tb = slide.shapes.add_textbox(px(a["x"]), px(a["y"]), px(a["w"]), px(a["h"]))
    _make_transparent(tb)
    _no_line(tb)

    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Pt(a.get("paddingLeft", 0) * 0.6)
    tf.margin_right  = Pt(a.get("paddingRight", 0) * 0.6)
    tf.margin_top    = Pt(a.get("paddingTop", 0) * 0.6)
    tf.margin_bottom = Pt(a.get("paddingBottom", 0) * 0.6)

    base_size  = a["fontSize"]
    base_color = a["color"] or [0, 0, 0]
    align      = _align(a.get("align","left"))
    is_upper   = a.get("uppercase", False)

    def _emit_run(p, text, *, bold=False, italic=False, color=None, size=None):
        if is_upper: text = text.upper()
        r = p.add_run()
        r.text = text
        r.font.name      = FONT
        r.font.size      = Pt((size or base_size) * 0.75)
        r.font.bold      = bold
        r.font.italic    = italic
        r.font.color.rgb = RGBColor(*(color or base_color))

    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after  = Pt(0)
    # line-height
    if a.get("lineHeight"):
        # OOXML utiliza spcPct (millonésimas) o lnSpc
        ratio = a["lineHeight"] / a["fontSize"] if a["fontSize"] else 1.2
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(ratio * 100000)))

    # ::before
    if a.get("before"):
        b = a["before"]
        _emit_run(p, b["text"],
                  bold=(b["fontWeight"] >= 600),
                  color=b["color"] or base_color,
                  size=b["fontSize"] or base_size)

    # runs
    first_para = True
    for run in a["runs"]:
        text = run["text"]
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            continue
        # split por newline interno
        parts = text.split("\n")
        for k, part in enumerate(parts):
            if k > 0:
                p = tf.add_paragraph()
                p.alignment = align
            if part:
                _emit_run(p, part,
                          bold=run["bold"],
                          italic=run["italic"],
                          color=run["color"] or base_color)

    # ::after
    if a.get("after"):
        af = a["after"]
        _emit_run(p, af["text"],
                  bold=(af["fontWeight"] >= 600),
                  color=af["color"] or base_color,
                  size=af["fontSize"] or base_size)


# ─── pipeline principal ───────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Emu(W_EMU)
    prs.slide_height = Emu(H_EMU)
    blank = prs.slide_layouts[6]

    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx = browser.new_context(viewport={"width": W_CSS, "height": H_CSS})
        page = ctx.new_page()
        page.goto(INDEX.as_uri(), wait_until="networkidle")
        page.evaluate("document.fonts.ready")
        page.emulate_media(media="print")
        page.wait_for_timeout(800)

        n = page.locator(".slide-container").count()
        print(f"[info] {n} slides", file=sys.stderr)

        for i in range(n):
            atoms = page.evaluate(JS_EXTRACT, i)
            slide = prs.slides.add_slide(blank)
            shapes_n = 0
            texts_n  = 0
            for a in atoms:
                try:
                    if a["kind"] == "shape":
                        add_shape_atom(slide, a)
                        shapes_n += 1
                    elif a["kind"] == "text":
                        add_text_atom(slide, a)
                        texts_n += 1
                except Exception as e:
                    print(f"[warn] slide {i+1} atom: {e}", file=sys.stderr)
            print(f"[pptx] {i+1:02d} shapes={shapes_n}  texts={texts_n}", file=sys.stderr)

        browser.close()

    prs.save(OUT_PPTX)
    print(f"[done] {OUT_PPTX.name} ({OUT_PPTX.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    build()
