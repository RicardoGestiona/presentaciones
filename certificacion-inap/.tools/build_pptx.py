"""
Genera PPTX 1280x720 desde index.html.

Estrategia: emulate_media('print') hace todas las slides visibles
verticalmente (CSS @media print del index.html). Captura cada
.slide-container con element.screenshot() y empaqueta en PPTX
con 1 imagen full-bleed por slide. Compatible Google Slides
(import nativo PPTX).
"""

from pathlib import Path
import sys

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
INDEX = ROOT / "index.html"
SHOTS_DIR = ROOT / ".tools" / "shots"
OUTPUT = ROOT / "certificacion-inap.pptx"

SLIDE_W_PX = 1280
SLIDE_H_PX = 720
EMU_PER_PX = 9525  # 96 DPI


def capture_slides() -> list[Path]:
    SHOTS_DIR.mkdir(parents=True, exist_ok=True)
    for old in SHOTS_DIR.glob("slide_*.png"):
        old.unlink()

    paths: list[Path] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        ctx = browser.new_context(
            viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
            device_scale_factor=2,
        )
        page = ctx.new_page()
        page.goto(INDEX.as_uri(), wait_until="networkidle")
        page.emulate_media(media="print")
        page.wait_for_timeout(500)

        slides = page.locator(".slide-container")
        total = slides.count()
        print(f"[info] {total} slides detectadas", file=sys.stderr)

        for i in range(total):
            target = SHOTS_DIR / f"slide_{i + 1:02d}.png"
            slides.nth(i).screenshot(path=str(target), omit_background=False)
            paths.append(target)
            print(f"[ok]   slide {i + 1:02d} -> {target.name}", file=sys.stderr)

        browser.close()
    return paths


def build_pptx(images: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_PX * EMU_PER_PX)
    prs.slide_height = Emu(SLIDE_H_PX * EMU_PER_PX)
    blank = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(OUTPUT)
    print(f"[done] {OUTPUT.relative_to(ROOT)} ({OUTPUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    imgs = capture_slides()
    if not imgs:
        sys.exit("error: no se generaron screenshots")
    build_pptx(imgs)
