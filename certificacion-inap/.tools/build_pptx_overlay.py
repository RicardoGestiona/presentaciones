"""
build_pptx_overlay.py — PPTX fiel + editable + PDF.

Estrategia:
  - Screenshot full-bleed como fondo  → fidelidad visual 100 %
  - Textboxes transparentes superpuestos → texto seleccionable y editable
  - Playwright extrae posición y estilo de cada elemento de texto
  - Módulos: además tabla nativa editable sobre la ficha
"""
from __future__ import annotations

import re, sys
from pathlib import Path
from typing import Optional

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

# ── rutas ──────────────────────────────────────────────────────────────────
ROOT      = Path(__file__).resolve().parent.parent
INDEX     = ROOT / "index.html"
SHOTS_DIR = ROOT / ".tools" / "shots"
OUT_PPTX  = ROOT / "certificacion-inap.pptx"
OUT_PDF   = ROOT / "certificacion-inap.pdf"

# ── dimensiones ────────────────────────────────────────────────────────────
W_CSS, H_CSS = 1280, 720      # CSS px (coordenadas de bounding_box)
SCALE        = 1              # 1:1 px→EMU sin ambigüedad de escala
EMU_PX       = 9525
W_EMU        = W_CSS * EMU_PX
H_EMU        = H_CSS * EMU_PX

def px(n: float) -> Emu:
    return Emu(int(n * EMU_PX))

# ── paleta ─────────────────────────────────────────────────────────────────
PRIMARY      = RGBColor(0,   109, 133)
PRIMARY_DARK = RGBColor(0,    77,  94)
WHITE        = RGBColor(255, 255, 255)
TEXT_MUTED   = RGBColor(180, 215, 225)
FONT         = "Roboto"

# ── JS: extrae todos los nodos de texto editables de un slide ──────────────
JS_GET_ELEMENTS = """
(slideIndex) => {
    const slides = document.querySelectorAll('.slide-container');
    const slide  = slides[slideIndex];
    if (!slide) return [];
    const sr = slide.getBoundingClientRect();

    const SKIP = new Set(['SCRIPT','STYLE','SVG','PATH','IMG']);
    const results = [];

    function walk(el) {
        if (SKIP.has(el.tagName)) return;
        // skip slide-number
        if (el.classList && el.classList.contains('slide-number')) return;

        const children = Array.from(el.children);
        const hasBlockChild = children.some(c => {
            const d = window.getComputedStyle(c).display;
            return d === 'block' || d === 'flex' || d === 'grid' || d === 'table';
        });

        // si tiene hijos bloque, bajar; si es hoja o solo inline, capturar
        if (!hasBlockChild && el.innerText && el.innerText.trim()) {
            const r  = el.getBoundingClientRect();
            const st = window.getComputedStyle(el);
            if (r.width > 2 && r.height > 2) {
                const rgb = st.color.match(/\\d+/g) || [0,0,0];
                let align = st.textAlign;
                if (align === 'start') align = 'left';
                if (align === 'end')   align = 'right';
                results.push({
                    text:       el.innerText.trim(),
                    x:          r.left   - sr.left,
                    y:          r.top    - sr.top,
                    w:          r.width,
                    h:          r.height,
                    fontSize:   parseFloat(st.fontSize),
                    fontWeight: parseInt(st.fontWeight) || 400,
                    r: parseInt(rgb[0]), g: parseInt(rgb[1]), b: parseInt(rgb[2]),
                    align:      align,
                    italic:     st.fontStyle === 'italic',
                });
            }
        } else {
            children.forEach(walk);
        }
    }
    walk(slide);
    return results;
}
"""

# ── JS: detecta si slide es módulo ────────────────────────────────────────
JS_IS_MODULE = """
(slideIndex) => {
    const slides = document.querySelectorAll('.slide-container');
    const slide  = slides[slideIndex];
    const h2 = slide && slide.querySelector('h2.slide-title');
    return h2 ? /^Módulo \\d+/.test(h2.textContent.trim()) : false;
}
"""

# ── JS: bounding box de la ficha del módulo ───────────────────────────────
JS_FICHA_BBOX = """
(slideIndex) => {
    const slides = document.querySelectorAll('.slide-container');
    const slide  = slides[slideIndex];
    const sr     = slide.getBoundingClientRect();
    const ficha  = slide.querySelector("div[style*='background: var(--color-primary)']");
    if (!ficha) return null;
    const r = ficha.getBoundingClientRect();
    return { x: r.left-sr.left, y: r.top-sr.top, w: r.width, h: r.height };
}
"""

# ── playwright: captura todo ───────────────────────────────────────────────

def capture_all() -> tuple[list[dict], list[Path]]:
    """
    Devuelve:
      slides_meta: [{is_module, elements, ficha_box}]
      shot_paths:  [Path a PNG por slide]
    """
    SHOTS_DIR.mkdir(parents=True, exist_ok=True)
    slides_meta: list[dict] = []
    shot_paths:  list[Path] = []

    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx = browser.new_context(
            viewport={"width": W_CSS, "height": H_CSS},
            device_scale_factor=SCALE,
        )
        page = ctx.new_page()
        page.goto(INDEX.as_uri(), wait_until="networkidle")
        page.emulate_media(media="print")
        page.wait_for_timeout(600)

        containers = page.locator(".slide-container")
        total      = containers.count()
        print(f"[info] {total} slides", file=sys.stderr)

        for i in range(total):
            # screenshot con texto invisible → solo shapes/colores/backgrounds
            page.evaluate("""
                document.querySelectorAll(
                  '.slide-container *:not(div):not(section), .slide-container span, .slide-container strong'
                ).forEach(el => {
                    const st = window.getComputedStyle(el);
                    if (st.display !== 'none' && el.childElementCount === 0) {
                        el.dataset._origColor = el.style.color;
                        el.style.setProperty('color', 'transparent', 'important');
                    }
                });
                // ocultar textos en divs hoja también
                document.querySelectorAll('.slide-container div, .slide-container p, .slide-container h1, .slide-container h2, .slide-container h3, .slide-container h4, .slide-container li, .slide-container td')
                .forEach(el => {
                    if (el.childElementCount === 0 && el.innerText.trim()) {
                        el.dataset._origColor = el.style.color;
                        el.style.setProperty('color', 'transparent', 'important');
                    }
                });
            """)
            target = SHOTS_DIR / f"slide_{i+1:02d}.png"
            containers.nth(i).screenshot(path=str(target))
            # restaurar colores
            page.evaluate("""
                document.querySelectorAll('[data-orig-color]').forEach(el => {
                    el.style.color = el.dataset._origColor || '';
                    delete el.dataset._origColor;
                });
                // restaurar todo lo que pusimos transparent
                document.querySelectorAll('.slide-container *').forEach(el => {
                    if (el.style.color === 'transparent') el.style.color = '';
                });
            """)
            page.wait_for_timeout(50)
            shot_paths.append(target)

            # metadatos de texto
            elements  = page.evaluate(JS_GET_ELEMENTS, i)
            is_module = page.evaluate(JS_IS_MODULE, i)
            ficha_box = page.evaluate(JS_FICHA_BBOX, i) if is_module else None

            slides_meta.append({
                "is_module": is_module,
                "elements":  elements,
                "ficha_box": ficha_box,
            })
            print(f"[info] slide {i+1:02d}  module={is_module}  elems={len(elements)}", file=sys.stderr)

        browser.close()
    return slides_meta, shot_paths

# ── helpers PPTX ──────────────────────────────────────────────────────────

def _align(a: str) -> PP_ALIGN:
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}.get(a, PP_ALIGN.LEFT)

def _make_transparent(shape) -> None:
    """Fuerza noFill real en XML — fill.background() no es transparente."""
    sp   = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    # eliminar cualquier fill existente
    for tag in ("a:solidFill","a:gradFill","a:pattFill","a:blipFill","a:grpFill","a:noFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    etree.SubElement(spPr, qn("a:noFill"))
    # sin línea
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for ch in list(ln): ln.remove(ch)
    etree.SubElement(ln, qn("a:noFill"))


def add_text_overlay(slide, el: dict) -> None:
    """Textbox con noFill real superpuesto exactamente sobre el elemento."""
    tb = slide.shapes.add_textbox(px(el["x"]), px(el["y"]), px(el["w"]), px(el["h"]))
    _make_transparent(tb)
    tf = tb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = _align(el.get("align","left"))
    run = p.add_run()
    run.text           = el["text"]
    run.font.name      = FONT
    run.font.size      = Pt(el["fontSize"] * 0.75)   # px → pt
    run.font.bold      = el.get("fontWeight", 400) >= 600
    run.font.italic    = el.get("italic", False)
    run.font.color.rgb = RGBColor(el["r"], el["g"], el["b"])

# ── tabla ficha editable ───────────────────────────────────────────────────

def _cell_solid(cell, color: RGBColor) -> None:
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")): tcPr.remove(old)
    sf   = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(sf,   qn("a:srgbClr"))
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")

def _cell_no_border(cell) -> None:
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for b in ("a:lnL","a:lnR","a:lnT","a:lnB"):
        el = tcPr.find(qn(b))
        if el is None: el = etree.SubElement(tcPr, qn(b))
        for ch in list(el): el.remove(ch)
        etree.SubElement(el, qn("a:noFill"))

def _cell_write(cell, text, pt, bold=False, color=None, margin_top=8) -> None:
    tf = cell.text_frame
    tf.clear()
    tf.margin_top    = Pt(margin_top)
    tf.margin_bottom = Pt(2)
    tf.margin_left   = Pt(8)
    tf.margin_right  = Pt(4)
    tf.word_wrap     = True
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.color.rgb = color or WHITE

def add_ficha_table(slide, box: dict) -> None:
    x, y, w, h = box["x"], box["y"], box["w"], box["h"]
    ficha = {"HORAS DE FORMACIÓN":"—","SESIONES ONLINE":"—",
             "SESIONES PRESENCIALES":"—","CLAUSTRO":"—"}
    rows  = list(ficha.items())
    n     = len(rows)

    HEADER_H = 52
    TABLE_Y  = y + HEADER_H
    TABLE_H  = h - HEADER_H - 4
    ROW_H    = TABLE_H // n

    t = slide.shapes.add_table(n, 2, px(x), px(TABLE_Y), px(w), px(TABLE_H)).table
    t.columns[0].width = px(w * 0.56)
    t.columns[1].width = px(w * 0.44)

    for i, (label, val) in enumerate(rows):
        row        = t.rows[i]
        row.height = px(ROW_H)
        lc, vc     = row.cells[0], row.cells[1]
        _cell_solid(lc, PRIMARY_DARK)
        _cell_solid(vc, PRIMARY)
        _cell_no_border(lc)
        _cell_no_border(vc)
        _cell_write(lc, label, 7.5, bold=True, color=TEXT_MUTED)
        _cell_write(vc, val,   15,  bold=True, color=WHITE)

# ── construir PPTX ─────────────────────────────────────────────────────────

def build_pptx(slides_meta: list[dict], shot_paths: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width  = Emu(W_EMU)
    prs.slide_height = Emu(H_EMU)
    blank = prs.slide_layouts[6]

    for i, (meta, img) in enumerate(zip(slides_meta, shot_paths)):
        slide = prs.slides.add_slide(blank)

        # 1. fondo: screenshot (fidelidad visual 100 %)
        if img.exists():
            slide.shapes.add_picture(str(img), 0, 0, Emu(W_EMU), Emu(H_EMU))

        # 2. overlays de texto transparentes
        for el in meta["elements"]:
            try:
                add_text_overlay(slide, el)
            except Exception as e:
                print(f"[warn] slide {i+1} elem skip: {e}", file=sys.stderr)

        # 3. módulos: tabla ficha editable
        if meta["is_module"] and meta.get("ficha_box"):
            add_ficha_table(slide, meta["ficha_box"])

        lbl = "módulo" if meta["is_module"] else "slide"
        print(f"[pptx] {i+1:02d} {lbl} — {len(meta['elements'])} overlays", file=sys.stderr)

    prs.save(OUT_PPTX)
    print(f"[done] {OUT_PPTX.name} ({OUT_PPTX.stat().st_size // 1024} KB)")

# ── PDF ────────────────────────────────────────────────────────────────────

def export_pdf() -> None:
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx  = browser.new_context(viewport={"width": W_CSS, "height": H_CSS})
        page = ctx.new_page()
        page.goto(INDEX.as_uri(), wait_until="networkidle")
        page.emulate_media(media="print")
        page.wait_for_timeout(800)
        page.pdf(path=str(OUT_PDF), width=f"{W_CSS}px", height=f"{H_CSS}px",
                 print_background=True,
                 margin={"top":"0","right":"0","bottom":"0","left":"0"})
        browser.close()
    print(f"[done] {OUT_PDF.name} ({OUT_PDF.stat().st_size // 1024} KB)")

# ── main ───────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("[browser] capturando slides y extrayendo texto …", file=sys.stderr)
    slides_meta, shot_paths = capture_all()

    print("[pptx] construyendo …", file=sys.stderr)
    build_pptx(slides_meta, shot_paths)

    print("[pdf] generando …", file=sys.stderr)
    export_pdf()
